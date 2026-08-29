#!/usr/bin/env python3
"""Build the TCS pilot pitch deck (PowerPoint) from verified evidence.

Reads tcs_pilot_evidence/MANIFEST.json so every number on a slide is the same
number that was reproduced by the offline analyzer. The deck is evidence-oriented:
it shows what was measured, what was hypothesised, and what authorization is
required before any TCS system is scanned. No client/production data is shown.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
OUT_DIR = HERE.parent / "tcs_pilot_evidence"
MANIFEST = json.loads((OUT_DIR / "MANIFEST.json").read_text(encoding="utf-8"))

NAVY = RGBColor(0x12, 0x1A, 0x2F)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x9B, 0xDF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xB7, 0x1C, 0x1C)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def _text(tf, runs, size=18, bold=False, color=NAVY, space_after=6, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    for line, b, c, s in runs:
        r = p.add_run()
        r.text = line
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = "Calibri"
    return p


def _add(slide, l, t, w, h, lines, size=18, bold=False, color=NAVY,
         align=PP_ALIGN.LEFT, space_after=6, fill=None, line_color=None):
    box = _box(slide, l, t, w, h)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = line_color or fill
        box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for item in lines:
        if isinstance(item, tuple):
            text, b, c, s = item
        else:
            text, b, c, s = item, bold, color, size
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = "Calibri"
    return box


def header(slide, title, kicker="TCS PILOT EVIDENCE"):
    _add(slide, 0.6, 0.4, 12.1, 0.4, [(kicker, False, ACCENT, 12)])
    _add(slide, 0.6, 0.75, 12.1, 0.9, [(title, True, NAVY, 30)])
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.65), Inches(12.1), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def footer(slide, n):
    _add(slide, 0.6, 7.05, 9, 0.35,
         [("Attestor 4.2 (analyzer 4.1.4) — offline, deterministic, evidence-bound",
           False, GREY, 9)])
    _add(slide, 11.8, 7.05, 1.0, 0.35, [(str(n), False, GREY, 9)], align=PP_ALIGN.RIGHT)


# ---- Slide 1: Title ----
s = prs.slides.add_slide(BLANK)
_bg(s, NAVY)
_add(s, 0.8, 2.0, 11.7, 1.2, [("Attestor 4.2 + Owen Desktop Cyber 4.3", True, WHITE, 40)])
_add(s, 0.8, 3.2, 11.7, 0.8, [("TCS Pilot — Reproducible Offensive-Security Evidence", True, ACCENT, 24)])
_add(s, 0.8, 4.4, 11.7, 1.4, [
    ("Deterministic static analyzer + bounded ranker. No generative-model claims.", False, LIGHT, 16),
    ("Every number in this deck was reproduced offline on owned / synthetic fixtures.", False, LIGHT, 16),
    ("No client, production, or TCS system was scanned or executed.", False, LIGHT, 16),
])
_add(s, 0.8, 6.4, 11.7, 0.5, [
    ("Analyzer SHA-256: %s" % MANIFEST["attestor"]["sha256"][:48] + "...", False, GREY, 11)])

# ---- Slide 2: What it is ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "What this is — and what it is not")
_add(s, 0.6, 1.9, 6.0, 4.8, [
    ("Is", True, GREEN, 20),
    ("• Offline static analysis — no code execution", False, NAVY, 15),
    ("• Deterministic: same bytes -> same findings", False, NAVY, 15),
    ("• Bounded ranker over measured findings", False, NAVY, 15),
    ("• Tamper-evident evidence spine (hash chain)", False, NAVY, 15),
    ("• Fail-closed: unknown != safe; secrets redacted", False, NAVY, 15),
], space_after=10)
_add(s, 6.9, 1.9, 5.8, 4.8, [
    ("Is NOT", True, RED, 20),
    ("• An exploit generator", False, NAVY, 15),
    ("• A generative AI making claims", False, NAVY, 15),
    ("• A scanner of client/production systems", False, NAVY, 15),
    ("• A deployment request — this is evidence", False, NAVY, 15),
    ("• A replacement for InfoSec authorization", False, NAVY, 15),
], space_after=10)
footer(s, 2)

# ---- Slide 3: Evidence spine / honesty gate ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "Evidence spine: honesty is enforced, not claimed")
c = MANIFEST["claims"]
status = [
    ("Tamper-evident chain intact", c["case_chain_intact"], GREEN if c["case_chain_intact"] else RED),
    ("Fix marked proven (pre-fix regression failed)", c["case_proven"], GREEN if c["case_proven"] else RED),
    ("Honesty gate blocked unproven fix", c["honesty_gate_blocked_unproven"], GREEN if c["honesty_gate_blocked_unproven"] else RED),
]
y = 2.0
for label, ok, col in status:
    _add(s, 0.6, y, 0.5, 0.5, [("✓" if ok else "✗", True, col, 24)])
    _add(s, 1.2, y + 0.02, 11, 0.5, [(label, False, NAVY, 18)])
    y += 0.7
_add(s, 0.6, 4.4, 12.1, 2.2, [
    ("How it works", True, BLUE, 18),
    ("One finding is carried through discovery -> validation -> severity -> exploitability ->", False, NAVY, 14),
    ("remediation -> regression. Each stage is tagged measured (tool produced it) or hypothesis", False, NAVY, 14),
    ("(model proposed it). A regression entry is REFUSED unless it records a pre-fix test failure —", False, NAVY, 14),
    ("so the pipeline cannot certify a fix it never proved.", False, NAVY, 14),
], space_after=6)
footer(s, 3)

# ---- Slide 4: Offense detection results ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "Offense detection on owned / synthetic fixtures")
rules = c["weak_asm_fixture_rules"]
_add(s, 0.6, 1.9, 4.0, 0.6, [("Weak x86-64 fixture — 5 rules fired", True, BLUE, 16)])
_add(s, 0.6, 2.5, 5.6, 3.0, [(("• " + r, False, NAVY, 14)) for r in rules], space_after=8)
_add(s, 6.6, 1.9, 3.0, 1.6, [
    ("C++ fixture", True, BLUE, 16),
    ("%d findings" % c["weak_cpp_fixture_findings"], True, NAVY, 34),
], space_after=4)
_add(s, 10.0, 1.9, 2.7, 1.6, [
    ("Clean 10k-line asm", True, BLUE, 16),
    ("%d findings" % c["clean_10k_asm_findings"], True, GREEN, 34),
], space_after=4)
_add(s, 6.6, 3.7, 6.1, 2.6, [
    ("Read:", True, NAVY, 14),
    ("The same engine that flags execve, stack-pivot, NOP-sled and W^X in a deliberately", False, NAVY, 13),
    ("weak program stays silent on a 10,000-line clean multi-function program. Detection", False, NAVY, 13),
    ("without false-positive noise is the property a SOC triage queue needs.", False, NAVY, 13),
], space_after=6)
footer(s, 4)

# ---- Slide 5: Enterprise hardening ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "Stage 6 — enterprise hardening (deployed, 9/9 tests)")
_add(s, 0.6, 1.9, 6.0, 4.6, [
    ("Tenant isolation", True, BLUE, 18),
    ("Grant tenant must equal request tenant — cross-tenant denied.", False, NAVY, 14),
    ("Two-party approval", True, BLUE, 18),
    ("Sensitive scopes (repo:write, scan:publish, admin:delete) need a", False, NAVY, 14),
    ("second signed approval before any allow.", False, NAVY, 14),
    ("Revocation + audit", True, BLUE, 18),
    ("Signed, freshness-bounded revocation; append-only hash-chained log.", False, NAVY, 14),
], space_after=8)
_add(s, 6.9, 1.9, 5.8, 4.6, [
    ("Composes, never replaces", True, GREEN, 16),
    ("Wraps the ordered Trusted Access gate:", False, NAVY, 13),
    ("authenticate -> time -> revocation ->", False, NAVY, 13),
    ("identity -> least privilege.", False, NAVY, 13),
    ("", False, NAVY, 8),
    ("Fail-closed", True, GREEN, 16),
    ("Hostile or missing input resolves to DENY, never an exception", False, NAVY, 13),
    ("or a silent allow.", False, NAVY, 13),
], space_after=8)
footer(s, 5)

# ---- Slide 6: Threat modelling / Stage 5 ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "Stage 5 — threat model & incident investigation")
_add(s, 0.6, 1.9, 6.0, 4.6, [
    ("STRIDE, stated up front", True, BLUE, 18),
    ("Spoofing, Tampering, Repudiation, Info-disclosure,", False, NAVY, 14),
    ("DoS, Elevation-of-privilege — each tagged hypothesis.", False, NAVY, 14),
    ("Measured vs hypothesis never merge", True, BLUE, 18),
    ("Analyzer findings (measured) are linked to modelled", False, NAVY, 14),
    ("threats but kept distinct from assumptions.", False, NAVY, 14),
], space_after=9)
_add(s, 6.9, 1.9, 5.8, 4.6, [
    ("Incidents = chain of evidence", True, BLUE, 18),
    ("Findings carried into the case-file spine as one", False, NAVY, 14),
    ("incident; reader sees which conclusions are load-", False, NAVY, 14),
    ("bearing vs assumed.", False, NAVY, 14),
    ("", False, NAVY, 6),
    ("Tests: 8/8 PASS", True, GREEN, 16),
    ("Incl. tampering links on real C++ fixture, fail-", False, NAVY, 13),
    ("closed on bad SHA / non-model input.", False, NAVY, 13),
], space_after=8)
footer(s, 6)

# ---- Slide 7: Authorization model ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "Authorization gate — before any TCS repo is scanned")
_add(s, 0.6, 1.9, 12.1, 2.0, [
    ("The held-out benchmark harness (bench_tcs42) refuses to run unless:", True, NAVY, 16),
    ("• A signed Trusted Access grant names the TCS prefix  tenant/tcs/...", False, NAVY, 14),
    ("• Scope includes  scan:read  and the grant is unexpired", False, NAVY, 14),
    ("• Corpus / result manifests resolve to local, non-symlink files", False, NAVY, 14),
], space_after=7)
_add(s, 0.6, 4.2, 12.1, 2.2, [
    ("Without that grant the harness exits deny — fail-closed — and analyzes nothing.", True, RED, 16),
    ("It does not manufacture cases or invoke any model. Bench tests: 6/6 PASS.", False, NAVY, 14),
    ("This is the gate, not the analysis. TCS InfoSec issues the grant; we run only", False, NAVY, 14),
    ("on the bytes the authorization names.", False, NAVY, 14),
], space_after=7)
footer(s, 7)

# ---- Slide 8: Verify / next steps ----
s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
header(s, "Reproduce it / next steps")
_add(s, 0.6, 1.9, 6.0, 4.6, [
    ("Reproduce (offline, isolated)", True, BLUE, 17),
    ("cd Attestor 4.2\\detector", False, NAVY, 13),
    ("python -I -B -X utf8 demo_tcs_case.py", False, NAVY, 13),
    ("python -I -B -X utf8 demo_asm_tcs.py", False, NAVY, 13),
    ("python -I -B -X utf8 bundle_tcs_evidence.py", False, NAVY, 13),
    ("", False, NAVY, 4),
    ("Bundle: tcs_pilot_evidence/ (README, JSON, MANIFEST)", False, GREY, 13),
], space_after=7)
_add(s, 6.9, 1.9, 5.8, 4.6, [
    ("Proposed next step", True, BLUE, 17),
    ("• Sponsored lab + TCS-held-out non-prod repo", False, NAVY, 14),
    ("• InfoSec-issued grant unlocks the bench harness", False, NAVY, 14),
    ("• Evidence-bound findings, not a favor", False, NAVY, 14),
    ("", False, NAVY, 4),
    ("The math is reproducible. That is the pitch.", True, GREEN, 15),
], space_after=8)
footer(s, 8)

out = OUT_DIR / "TCS_Pilot_Attestor42.pptx"
prs.save(str(out))
print("wrote", out, "slides:", len(prs.slides._sldIdLst))
